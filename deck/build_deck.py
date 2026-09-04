# -*- coding: utf-8 -*-
"""
Builds the SIH 2026 Idea Submission deck for Problem Statement 26126 (BEL)
on top of the official SIH2026-IDEA-Presentation-Format.pptx template.

Template chrome is preserved; each content slide carries one flat,
colour-blocked diagram drawn from native PowerPoint shapes, so the team can
edit any element and the export stays vector-crisp.

Run prep_logos.py once first to populate assets/logos.
"""
import os
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
SRC = os.path.join(ROOT, "source", "SIH2026-IDEA-Presentation-Format.pptx")
OUT = os.path.join(HERE, "SIH2026_PS26126_Idea_Presentation.pptx")
LOGOS = os.path.join(HERE, "assets", "logos")

TEAM_NAME = "THE VIKINGS"
TEAM_ID = "<Team ID>"          # <-- fill from the SIH portal after registration
FINALIST = "SIH 2025 FINALIST"

# ---------------------------------------------------------------- palette
GREEN = "34B77C"
RED = "E8544F"
YELLOW = "C79A0B"
ORANGE = "DE8429"
BLUE = "4A86D8"
PURPLE = "8659D6"
PINK = "DC4A96"
TEAL = "23A8C4"
NAVY = "1F2A37"
GREY = "8B9099"
DARK = "343A40"
MUTED = "5A616B"
WHITE = "FFFFFF"
PAPER = "F4F6F8"
HAIR = "DDE1E6"
GOLD = "E9C46A"
GOLD_DK = "8A6410"
GOLD_PALE = "FBF1D9"
GOLD_LINE = "D8B84A"

FONT = "Segoe UI"

L, R = 0.45, 12.88          # content left / right edge
W = R - L                   # 12.43 in usable width
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------- helpers
def rgb(h):
    return RGBColor.from_string(h)


def flat(shape, fill=None, line=None, line_w=1.0):
    """Remove theme shadow, apply flat fill and optional outline."""
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_w)
    return shape


def no_bullet(paragraph):
    """Strip any inherited/explicit bullet so template list styles don't leak in."""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buFont", "a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    el = parse_xml('<a:buNone xmlns:a="%s"/>' % NS_A)
    pPr.insert_element_before(el, "a:tabLst", "a:defRPr", "a:extLst")


def write(tf, specs, anchor=MSO_ANCHOR.MIDDLE, wrap=True, m=0.05):
    """specs: list of (text, size, bold, colour, align, space_before_pt)"""
    tf.clear()
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(m)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, (text, size, bold, colour, align, sb) in enumerate(specs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        no_bullet(p)
        p.alignment = align
        if sb:
            p.space_before = Pt(sb)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = rgb(colour)
    return tf


def box(slide, x, y, w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE, fill=None,
        line=None, radius=0.14, line_w=1.0, text=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    flat(s, fill, line, line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if text:
        write(s.text_frame, text)
    return s


def label(slide, x, y, w, h, specs, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    write(tb.text_frame, specs, anchor=anchor)
    return tb


def section(slide, x, y, w, text, colour=NAVY, size=12.5):
    return label(slide, x, y, w, 0.32, [(text, size, True, colour, PP_ALIGN.LEFT, 0)])


def bullets(slide, x, y, w, h, rows, size=10, colour=DARK, gap=7):
    return label(slide, x, y, w, h,
                 [("•  " + r, size, False, colour, PP_ALIGN.LEFT, 0 if i == 0 else gap)
                  for i, r in enumerate(rows)], anchor=MSO_ANCHOR.TOP)


def line(slide, x1, y1, x2, y2, colour=GREY, w=1.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(colour)
    c.line.width = Pt(w)
    return c


def logo(slide, name, cx, cy, max_w, max_h):
    """Place a logo scaled to fit max_w x max_h, centred on (cx, cy)."""
    path = os.path.join(LOGOS, name + ".png")
    iw, ih = Image.open(path).size
    s = min(max_w / iw, max_h / ih)
    w, h = iw * s, ih * s
    return slide.shapes.add_picture(path, Inches(cx - w / 2), Inches(cy - h / 2),
                                    Inches(w), Inches(h))


def clear_instructions(slide):
    """Drop the template's grey instruction text boxes from a content slide."""
    for sh in list(slide.shapes):
        if sh.name.startswith("TextBox"):
            sh._element.getparent().remove(sh._element)


def set_badge(slide):
    """Team identity block: name plate + SIH 2025 finalist ribbon."""
    for sh in slide.shapes:
        if sh.has_text_frame and "Your Team Name" in sh.text_frame.text:
            sh.left, sh.top = Inches(0.30), Inches(0.14)
            sh.width, sh.height = Inches(2.00), Inches(0.66)
            flat(sh, fill=NAVY)
            write(sh.text_frame, [(TEAM_NAME, 12, True, WHITE, PP_ALIGN.CENTER, 0)], m=0.02)
    box(slide, 0.30, 0.86, 2.00, 0.30, fill=GOLD, radius=0.40,
        text=[(FINALIST, 8.5, True, GOLD_DK, PP_ALIGN.CENTER, 0)])


# ---------------------------------------------------------------- build
prs = Presentation(SRC)
s1, s2, s3, s4, s5, s6 = [prs.slides[i] for i in range(6)]

# ================================================================ SLIDE 1
for sh in s1.shapes:
    if sh.name == "TextBox 9":
        write(sh.text_frame, [
            ("Problem Statement ID  –  26126", 20, True, DARK, PP_ALIGN.LEFT, 0),
            ("Problem Statement Title  –  Vision Based Autonomous Navigation for "
             "Unmanned Ground Vehicle for Outdoor environment", 16, False, DARK, PP_ALIGN.LEFT, 12),
            ("Theme  –  Smart Automation", 16, False, DARK, PP_ALIGN.LEFT, 12),
            ("PS Category  –  Software", 16, False, DARK, PP_ALIGN.LEFT, 12),
            ("Team ID  –  " + TEAM_ID, 16, False, DARK, PP_ALIGN.LEFT, 12),
            ("Team Name  –  The Vikings", 16, True, DARK, PP_ALIGN.LEFT, 12),
        ], anchor=MSO_ANCHOR.TOP)

box(s1, 0.45, 5.62, 5.95, 0.86, fill=GOLD_PALE, line=GOLD_LINE, radius=0.16, line_w=1.5, text=[
    ("SIH 2025  GRAND FINALE FINALIST", 15, True, GOLD_DK, PP_ALIGN.CENTER, 0),
    ("Team The Vikings — returning to Smart India Hackathon 2026",
     10.5, False, MUTED, PP_ALIGN.CENTER, 3)])

# ================================================================ SLIDE 2
clear_instructions(s2)
set_badge(s2)

box(s2, L, 1.24, W, 0.66, fill=NAVY, radius=0.10, text=[
    ("DRISHTI-UGV  —  a vision-first autonomy module that drives an outdoor UGV "
     "from Point A to Point B with no GPS", 16, True, WHITE, PP_ALIGN.CENTER, 0)])

PW = 6.00
LX, RX = L, 6.88

box(s2, LX, 2.04, PW, 0.46, fill=RED, radius=0.30,
    text=[("THE PROBLEM", 13, True, WHITE, PP_ALIGN.CENTER, 0)])
box(s2, RX, 2.04, PW, 0.46, fill=GREEN, radius=0.30,
    text=[("OUR SOLUTION", 13, True, WHITE, PP_ALIGN.CENTER, 0)])

problems = [
    (RED, "01", "GPS is denied or unreliable",
     "Tree cover, terrain masking and jamming remove any trusted fix."),
    (ORANGE, "02", "Open ground is not safe ground",
     "A ditch, mud patch or slope still reads as free space."),
    (YELLOW, "03", "Obstacles appear after planning",
     "People, animals and vehicles cross the path mid-mission."),
    (PINK, "04", "Vision degrades without warning",
     "Glare, shadow and low light silently break a camera-only stack."),
]
solutions = [
    (GREEN, "A", "Visual SLAM replaces GPS",
     "RTAB-Map fuses stereo and IMU into a drift-bounded pose."),
    (TEAL, "B", "Terrain is scored, not just occupied",
     "Slope, roughness and height variance become real driving cost."),
    (BLUE, "C", "Semantics fused into the costmap",
     "YOLO marks mud, water, rock and people; unknown is priced high."),
    (PURPLE, "D", "Safety supervisor outside the AI",
     "Deterministic code halts on stale sensors or lost pose."),
]

ROWS = [2.62, 3.50, 4.38, 5.26]
for col_x, items in ((LX, problems), (RX, solutions)):
    for (colour, tag, head, sub), ry in zip(items, ROWS):
        box(s2, col_x, ry, 0.62, 0.62, fill=colour, radius=0.22,
            text=[(tag, 14, True, WHITE, PP_ALIGN.CENTER, 0)])
        label(s2, col_x + 0.78, ry - 0.05, PW - 0.82, 0.34,
              [(head, 12.5, True, colour, PP_ALIGN.LEFT, 0)], anchor=MSO_ANCHOR.BOTTOM)
        label(s2, col_x + 0.78, ry + 0.29, PW - 0.82, 0.34,
              [(sub, 10, False, MUTED, PP_ALIGN.LEFT, 0)], anchor=MSO_ANCHOR.TOP)

section(s2, L, 6.00, 5.0, "INNOVATION  &  UNIQUENESS")
uniq = [
    (PURPLE, "Deterministic safety, not learned",
     "The stop decision never depends on model confidence."),
    (TEAL, "Camera-first, LiDAR-free terrain",
     "Traversability from vision geometry plus semantics."),
    (BLUE, "Proven by simulation at scale",
     "Randomised Isaac Sim missions score it before hardware."),
]
uw = (W - 2 * 0.21) / 3
for i, (colour, head, sub) in enumerate(uniq):
    ux = L + i * (uw + 0.21)
    box(s2, ux, 6.28, uw, 0.60, fill=PAPER, line=colour, radius=0.16, line_w=1.25, text=[
        (head, 11, True, colour, PP_ALIGN.CENTER, 0),
        (sub, 9.5, False, MUTED, PP_ALIGN.CENTER, 1)])

# ================================================================ SLIDE 3
clear_instructions(s3)
set_badge(s3)

section(s3, L, 1.22, 9.0, "PROCESSING PIPELINE  —  camera in, safe velocity out")

pipe = [
    (TEAL, "Stereo Camera\n+ IMU", "RGB • depth • inertial"),
    (PURPLE, "Perception AI", "YOLO detect + segment"),
    (BLUE, "Visual SLAM", "RTAB-Map pose & map"),
    (GREEN, "Traversability\nCostmap", "slope • roughness • class"),
    (ORANGE, "Nav2 Planner", "global route"),
    (PINK, "MPPI Controller", "local trajectory"),
    (RED, "Safety Supervisor", "veto • slow • stop"),
]
nw, gap = 1.55, 0.26
py, ph = 1.56, 1.06
for i, (colour, head, sub) in enumerate(pipe):
    nx = L + i * (nw + gap)
    box(s3, nx, py, nw, ph, fill=colour, radius=0.16, text=[
        (head, 10.5, True, WHITE, PP_ALIGN.CENTER, 0),
        (sub, 8.5, False, WHITE, PP_ALIGN.CENTER, 2)])
    if i < len(pipe) - 1:
        box(s3, nx + nw + 0.03, py + ph / 2 - 0.10, 0.20, 0.20,
            shape=MSO_SHAPE.RIGHT_ARROW, fill=GREY)

label(s3, L, 2.66, W, 0.26,
      [("Output: /cmd_vel to the UGV base  •  every block is a ROS 2 node behind a fixed topic "
        "and TF contract, so the identical graph runs in simulation and on the real vehicle",
        9.5, False, MUTED, PP_ALIGN.CENTER, 0)])

section(s3, L, 3.00, 9.0, "LANGUAGES, FRAMEWORKS  &  TOOLS")
tools = [
    ("python", "Python", "perception, tooling"),
    ("cpp", "C++", "real-time nodes"),
    ("ros", "ROS 2 Jazzy", "middleware"),
    ("nvidia", "Isaac Sim", "simulation"),
    ("pytorch", "PyTorch", "model training"),
    ("opencv", "OpenCV", "vision operations"),
    ("docker", "Docker", "reproducible builds"),
    ("ubuntu", "Ubuntu 24.04", "target platform"),
]
tw = (W - 7 * 0.16) / 8
for i, (img, name, role) in enumerate(tools):
    tx = L + i * (tw + 0.16)
    box(s3, tx, 3.34, tw, 1.24, fill=PAPER, line=HAIR, radius=0.12)
    logo(s3, img, tx + tw / 2, 3.72, tw - 0.50, 0.52)
    label(s3, tx + 0.05, 4.02, tw - 0.10, 0.26,
          [(name, 11, True, NAVY, PP_ALIGN.CENTER, 0)])
    label(s3, tx + 0.05, 4.26, tw - 0.10, 0.24,
          [(role, 8.5, False, MUTED, PP_ALIGN.CENTER, 0)])

section(s3, L, 4.76, 10.0,
        "IMPLEMENTATION METHODOLOGY  —  simulation-first; every phase ships a runnable artefact")
phases = [("0", "Environment"), ("1", "Sim navigation"), ("2", "Visual SLAM"), ("3", "Terrain layer"),
          ("4", "Perception"), ("5", "Safety logic"), ("6", "1000+ missions"), ("7", "Hardware")]
shades = [NAVY, TEAL, BLUE, GREEN, PURPLE, RED, ORANGE, PINK]
pw2 = (W - 7 * 0.14) / 8
for i, (num, name) in enumerate(phases):
    px = L + i * (pw2 + 0.14)
    box(s3, px, 5.10, pw2, 0.66, fill=shades[i], radius=0.18, text=[
        ("PHASE " + num, 8, True, WHITE, PP_ALIGN.CENTER, 0),
        (name, 10, True, WHITE, PP_ALIGN.CENTER, 1)])

box(s3, L, 5.94, W, 0.84, fill=PAPER, radius=0.10, text=[
    ("Design principle: the neural network only answers “what am I looking at?”",
     11.5, True, NAVY, PP_ALIGN.CENTER, 0),
    ("Coordinate transforms, mapping, planning, control and the stop decision stay in "
     "deterministic, auditable code — so a perception failure degrades into a safe halt, "
     "never a collision.", 9.5, False, MUTED, PP_ALIGN.CENTER, 3)])

# ================================================================ SLIDE 4
clear_instructions(s4)
set_badge(s4)

cols = [
    (GREEN, "F", "FEASIBLE TODAY", [
        "Nav2, RTAB-Map and elevation_mapping_cupy are mature ROS 2 packages",
        "Integration and tuning — nothing has to be invented",
        "Simulation-first: development starts with zero hardware",
        "Whole loop runs on one RTX workstation (32 GB RAM, 16 GB VRAM)"]),
    (RED, "R", "KEY RISKS", [
        "TF, timestamp and calibration errors — the top failure cause",
        "SLAM drift on featureless or low-texture terrain",
        "A ditch or water surface misread as drivable ground",
        "Simulation-to-reality gap in lighting and texture"]),
    (BLUE, "M", "MITIGATION", [
        "Validate the TF tree and clock sync before any AI work",
        "Stereo + IMU fusion, speed limits, automatic re-localisation",
        "Unknown terrain is priced expensive, never free",
        "Domain randomisation over light, texture, noise and obstacles"]),
    (ORANGE, "V", "VIABILITY BEYOND SIH", [
        "Identical ROS 2 topic and TF contract, simulation and vehicle",
        "Only sensor driver and motor interface change on hardware",
        "Profiles down to a Jetson-class edge computer",
        "Reusable across BEL surveillance and logistics platforms"]),
]
cw4 = (W - 3 * 0.24) / 4
for i, (colour, letter, head, rows) in enumerate(cols):
    cx = L + i * (cw4 + 0.24)
    box(s4, cx, 1.70, cw4, 3.10, fill=colour, radius=0.10)
    box(s4, cx + cw4 / 2 - 0.35, 1.32, 0.70, 0.70, shape=MSO_SHAPE.OVAL, fill=colour,
        line=WHITE, line_w=2.5, text=[(letter, 16, True, WHITE, PP_ALIGN.CENTER, 0)])
    label(s4, cx + 0.16, 2.14, cw4 - 0.32, 0.34,
          [(head, 12.5, True, WHITE, PP_ALIGN.LEFT, 0)], anchor=MSO_ANCHOR.TOP)
    bullets(s4, cx + 0.16, 2.58, cw4 - 0.32, 2.10, rows, size=10, colour=WHITE, gap=8)

section(s4, L, 5.00, 11.0,
        "ACCEPTANCE GATES  —  what we will measure, instead of claiming “100% accuracy”")
gates = [(GREEN, "≥ 95%", "collision-free\nrandomised missions"),
         (TEAL, "< 2%", "localisation drift\nover distance travelled"),
         (BLUE, "< 200 ms", "emergency-stop\nsoftware response"),
         (PURPLE, "≤ 100 ms", "perception\nlatency"),
         (ORANGE, "≥ 20 Hz", "control loop\nrate"),
         (PINK, "≥ 97%", "goal completion\nrate")]
gw = (W - 5 * 0.18) / 6
for i, (colour, big, sub) in enumerate(gates):
    gx = L + i * (gw + 0.18)
    box(s4, gx, 5.34, gw, 1.08, fill=PAPER, line=colour, radius=0.14, line_w=1.25, text=[
        (big, 17, True, colour, PP_ALIGN.CENTER, 0),
        (sub, 9.5, False, MUTED, PP_ALIGN.CENTER, 3)])

label(s4, L, 6.48, W, 0.30,
      [("Every gate is scored automatically from rosbag2 logs against Isaac Sim ground truth, "
        "so the numbers are reproducible rather than asserted.",
        9.5, False, GREY, PP_ALIGN.CENTER, 0)])

# ================================================================ SLIDE 5
clear_instructions(s5)
set_badge(s5)

box(s5, 5.42, 2.46, 2.50, 2.50, shape=MSO_SHAPE.OVAL, fill=NAVY, text=[
    ("DRISHTI-UGV", 14, True, WHITE, PP_ALIGN.CENTER, 0),
    ("vision-first\nGPS-denied\nautonomy", 10.5, False, "C8CED6", PP_ALIGN.CENTER, 5)])

left_items = [
    (ORANGE, "DEFENCE  &  BEL PLATFORMS",
     "Reconnaissance, patrol and resupply in jammed or GPS-denied terrain."),
    (RED, "SEARCH  &  RESCUE",
     "Entry into collapsed or flooded zones where GPS and comms have failed."),
    (GREEN, "AGRICULTURE  &  SURVEY",
     "Crop-row navigation and field inspection without RTK base stations."),
]
right_items = [
    (BLUE, "COST  &  IMPORT SUBSTITUTION",
     "Cameras replace an expensive LiDAR; the stack is owned in-country."),
    (PURPLE, "SAFETY BY DESIGN",
     "A deterministic supervisor halts the vehicle before a failure becomes a collision."),
    (PINK, "ONE REUSABLE PLATFORM",
     "The same ROS 2 contract retargets to delivery, mining and inspection UGVs."),
]
CW = 4.70
ROW_Y = [1.52, 3.02, 4.52]
DOT_Y = [ry + 0.31 for ry in ROW_Y]
HUB_Y = 3.71

# elbow leader lines: dot -> vertical bus -> hub
for bus, stub_from, stub_to, dot_edge in ((5.28, 5.28, 5.42, 5.15),    # left
                                          (8.05, 7.92, 8.05, 8.18)):   # right
    for dy in DOT_Y:
        line(s5, dot_edge, dy, bus, dy)
    line(s5, bus, DOT_Y[0], bus, DOT_Y[-1])
    line(s5, stub_from, HUB_Y, stub_to, HUB_Y)

for items, cx, align in ((left_items, L, PP_ALIGN.RIGHT), (right_items, 8.18, PP_ALIGN.LEFT)):
    for (colour, head, sub), ry in zip(items, ROW_Y):
        if align == PP_ALIGN.RIGHT:
            dot_x, bx = cx + CW - 0.34, cx
        else:
            dot_x, bx = cx, cx + 0.44
        box(s5, dot_x, ry + 0.14, 0.34, 0.34, shape=MSO_SHAPE.OVAL, fill=colour)
        label(s5, bx, ry, CW - 0.44, 0.36, [(head, 12, True, colour, align, 0)])
        label(s5, bx, ry + 0.38, CW - 0.44, 0.66,
              [(sub, 10, False, MUTED, align, 0)], anchor=MSO_ANCHOR.TOP)

box(s5, L, 5.96, W, 0.84, fill=PAPER, radius=0.10, text=[
    ("Measured, not claimed.", 11.5, True, NAVY, PP_ALIGN.CENTER, 0),
    ("Every benefit above is backed by a mission-level number — collision-free completion rate, "
     "localisation error against simulator ground truth, path efficiency, recovery rate and "
     "emergency-stop latency.", 9.5, False, MUTED, PP_ALIGN.CENTER, 3)])

# ================================================================ SLIDE 6
clear_instructions(s6)
set_badge(s6)

groups = [
    (BLUE, "NAVIGATION  &  CONTROL", [
        "ros-navigation / navigation2 — planners, costmaps, behaviour trees",
        "Nav2 MPPI controller — sampling-based predictive local control",
        "sacrover / traversability-nav2 — outdoor traversability costmap"]),
    (GREEN, "VISUAL LOCALISATION  &  SLAM", [
        "introlab / rtabmap_ros — ROS 2 stereo and RGB-D SLAM with Nav2",
        "UZ-SLAMLab / ORB_SLAM3 — visual-inertial benchmark (GPLv3)",
        "LARIAD / Offroad-Nav — closest end-to-end off-road reference"]),
    (ORANGE, "TERRAIN  &  PERCEPTION", [
        "leggedrobotics / elevation_mapping_cupy — GPU elevation layers",
        "ultralytics — detection and segmentation models (licence reviewed)",
        "DepthAnything / Depth-Anything-V2 — monocular depth fallback"]),
]
gw6 = (W - 2 * 0.24) / 3
for i, (colour, head, rows) in enumerate(groups):
    gx = L + i * (gw6 + 0.24)
    box(s6, gx, 1.42, gw6, 1.86, fill=PAPER, line=HAIR, radius=0.10)
    box(s6, gx, 1.42, gw6, 0.40, fill=colour, radius=0.28,
        text=[(head, 11, True, WHITE, PP_ALIGN.CENTER, 0)])
    bullets(s6, gx + 0.12, 1.90, gw6 - 0.24, 1.32, rows, size=9.5, gap=7)

section(s6, L, 3.44, 9.0, "SOURCE LINKS  &  OFFICIAL DOCUMENTATION")
links_left = [
    "Nav2 framework  —  github.com/ros-navigation/navigation2",
    "RTAB-Map ROS 2  —  github.com/introlab/rtabmap_ros",
    "Elevation Mapping CuPy  —  github.com/leggedrobotics/elevation_mapping_cupy",
    "Traversability Nav2  —  github.com/sacrover/traversability-nav2",
    "Offroad-Nav reference  —  github.com/LARIAD/Offroad-Nav",
]
links_right = [
    "Isaac Sim ROS 2 bridge  —  docs.isaacsim.omniverse.nvidia.com",
    "Isaac Sim workspaces  —  github.com/isaac-sim/IsaacSim-ros_workspaces",
    "Ultralytics YOLO  —  github.com/ultralytics/ultralytics",
    "Depth Anything V2  —  github.com/DepthAnything/Depth-Anything-V2",
    "Open3D  —  github.com/isl-org/Open3D",
]
for col, cx in ((links_left, L), (links_right, 6.88)):
    box(s6, cx, 3.78, 6.00, 1.72, fill=WHITE, line=HAIR, radius=0.08)
    label(s6, cx + 0.16, 3.88, 5.68, 1.56,
          [("›  " + t, 9.5, False, DARK, PP_ALIGN.LEFT, 0 if k == 0 else 9)
           for k, t in enumerate(col)], anchor=MSO_ANCHOR.TOP)

lic = [
    (GREEN, "MIT / BSD / Apache-2.0",
     "Nav2, RTAB-Map, elevation_mapping_cupy, Offroad-Nav, Open3D"),
    (YELLOW, "Licence-reviewed before use",
     "Ultralytics and Depth Anything V2 Small — terms checked"),
    (RED, "Excluded from the shipped build",
     "ORB-SLAM3 (GPLv3) used only as an offline benchmark"),
]
lw = (W - 2 * 0.20) / 3
for i, (colour, head, sub) in enumerate(lic):
    lx = L + i * (lw + 0.20)
    box(s6, lx, 5.66, lw, 0.78, fill=PAPER, line=colour, radius=0.14, line_w=1.25, text=[
        (head, 10.5, True, colour, PP_ALIGN.CENTER, 0),
        (sub, 9, False, MUTED, PP_ALIGN.CENTER, 3)])

label(s6, L, 6.50, W, 0.30,
      [("Licence position screened per component before integration; every third-party notice is "
        "preserved in the delivered repository.", 9, False, GREY, PP_ALIGN.CENTER, 0)])

# ---------------------------------------------------------------- trim slide 7
ids = prs.slides._sldIdLst
ids.remove(list(ids)[6])

prs.save(OUT)
print("written:", os.path.basename(OUT), "| slides:", len(prs.slides._sldIdLst))
